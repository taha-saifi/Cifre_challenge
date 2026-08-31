#!/usr/bin/env python3
"""Render deliverables/deck.md to an 8-slide .pptx.

deck.md is written as titre-message / corps / à dire. Only the title and the body reach
the slide; the "À dire" blocks are speaker material and go into the slide notes, which is
where they belong -- putting them on the slide would break the 8-slide format by volume
even while respecting it by count.
"""
from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Cm, Pt

ROOT = Path(__file__).resolve().parent.parent
SRC = ROOT / "deliverables" / "deck.md"
OUT = ROOT / "deliverables" / "deck.pptx"

ACCENT = RGBColor(0x0E, 0x74, 0x90)
INK = RGBColor(0x17, 0x1B, 0x21)
MUTED = RGBColor(0x5A, 0x64, 0x72)


def clean(text: str) -> str:
    """Strip markdown emphasis and code ticks -- pptx carries no inline markup here."""
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    text = re.sub(r"`([^`]+)`", r"\1", text)
    return re.sub(r"(?<!\*)\*([^*]+)\*(?!\*)", r"\1", text).strip()


def parse_slides(markdown: str) -> list[dict]:
    slides = []
    for block in re.split(r"\n## Slide ", markdown)[1:]:
        header, _, rest = block.partition("\n")
        slide = {"heading": header.strip(), "title": "", "body": [], "notes": []}
        section = None
        for line in rest.split("\n"):
            stripped = line.strip()
            if stripped.startswith("### "):
                label = stripped[4:].lower()
                section = "title" if "titre" in label else ("notes" if "dire" in label else "body")
                continue
            if not stripped or stripped.startswith("---") or stripped.startswith("## "):
                continue
            if section == "title" and not slide["title"]:
                slide["title"] = clean(stripped)
            elif section == "body":
                slide["body"].append(stripped)
            elif section == "notes":
                slide["notes"].append(clean(stripped))
        slides.append(slide)
    return slides


def body_lines(raw: list[str]) -> list[tuple[int, str]]:
    """Flatten the markdown body into (indent_level, text), keeping tables readable."""
    out: list[tuple[int, str]] = []
    for line in raw:
        if set(line) <= set("|-: "):
            continue
        if line.startswith("|"):
            cells = [clean(c) for c in line.strip().strip("|").split("|")]
            out.append((1, "  ·  ".join(c for c in cells if c)))
        elif re.match(r"^[-*] ", line):
            out.append((0, clean(line[2:])))
        elif re.match(r"^\d+\. ", line):
            out.append((0, clean(re.sub(r"^\d+\.\s*", "", line))))
        else:
            out.append((1, clean(line)))
    return out


def main() -> int:
    slides_data = parse_slides(SRC.read_text(encoding="utf-8"))
    if len(slides_data) != 8:
        print(f"ATTENTION : {len(slides_data)} slides trouvees, 8 attendues (§18.2)")

    presentation = Presentation()
    presentation.slide_width, presentation.slide_height = Cm(33.87), Cm(19.05)  # 16:9
    blank = presentation.slide_layouts[6]

    for index, data in enumerate(slides_data, start=1):
        slide = presentation.slides.add_slide(blank)

        eyebrow = slide.shapes.add_textbox(Cm(1.6), Cm(0.9), Cm(30), Cm(0.9)).text_frame
        eyebrow.text = f"{index:02d} — {data['heading']}"
        run = eyebrow.paragraphs[0].runs[0]
        run.font.size, run.font.bold, run.font.color.rgb = Pt(11), True, ACCENT

        title = slide.shapes.add_textbox(Cm(1.6), Cm(1.8), Cm(30.6), Cm(2.4)).text_frame
        title.word_wrap = True
        title.text = data["title"]
        run = title.paragraphs[0].runs[0]
        run.font.size, run.font.bold, run.font.color.rgb = Pt(26), True, INK

        lines = body_lines(data["body"])
        box = slide.shapes.add_textbox(Cm(1.6), Cm(4.6), Cm(30.6), Cm(13.2)).text_frame
        box.word_wrap = True
        # Slides carry a readable subset; the rest is speaker material in the notes.
        size = Pt(15) if len(lines) <= 9 else (Pt(13) if len(lines) <= 13 else Pt(11))
        for position, (level, text) in enumerate(lines):
            para = box.paragraphs[0] if position == 0 else box.add_paragraph()
            para.text = ("• " + text) if level == 0 else text
            para.level = level
            para.space_after = Pt(5)
            for run in para.runs:
                run.font.size = size
                run.font.color.rgb = INK if level == 0 else MUTED

        if data["notes"]:
            slide.notes_slide.notes_text_frame.text = " ".join(data["notes"])

    presentation.save(OUT)
    print(f"wrote {OUT} — {len(presentation.slides.__iter__.__self__._sldIdLst)} slides")
    for index, data in enumerate(slides_data, start=1):
        print(f"  {index}. {data['title'][:66]}  ({len(body_lines(data['body']))} lignes)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
